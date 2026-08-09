"""PNPI / PNPI · Endpoints d'exports (CSV, PDF)."""

from __future__ import annotations

import io
import os
import zipfile
from datetime import UTC, date

import qrcode
from fastapi import APIRouter, Body, Depends, HTTPException, Query
from fastapi.responses import Response, StreamingResponse
from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sqlalchemy import func, select
from sqlalchemy.orm import Session

from ..core.audit import write_audit_event
from ..core.auth import Role, User, require_roles
from ..core.executive_report import generate_executive_report
from ..core.signature import generate_signature
from ..core.storage import get_storage
from ..database import get_db, now_utc
from ..models.core import FieldReportORM, TraceBatchORM, UnitORM, UserAccountORM
from ..models.pilotage import ProjectDossierTransitionORM
from ..models.pnpi import (
    AgrementTechniqueIndustrielORM,
    DocumentDossierORM,
    InspectionConformiteORM,
    OperateurIndustrielORM,
)

router = APIRouter(tags=["Exports"])


@router.get("/exports/executive-report.pdf", summary="Rapport executif hebdomadaire")
async def export_executive_report(
    current_user: User = Depends(require_roles(Role.admin, Role.ministre, Role.directeur)),
    db: Session = Depends(get_db),
):
    pdf_bytes = generate_executive_report(db)
    now = now_utc()
    write_audit_event(
        db,
        actor=current_user.username,
        action="export.executive_report",
        target="weekly",
        details="Rapport executif hebdomadaire genere",
    )
    db.commit()
    return Response(
        content=pdf_bytes,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="rapport_executif_pnpi_{now.strftime("%Y%m%d")}.pdf"'},
    )


def _csv_generator(rows, header):
    """Stream CSV rows one at a time for large datasets."""
    yield ",".join(header) + "\n"
    for row in rows:
        yield ",".join(str(v) for v in row) + "\n"


def _filter_pilotage_transitions(rows, *, dossier_id, changed_by, date_from, date_to):
    filtered = list(rows)
    if dossier_id:
        normalized = dossier_id.strip().upper()
        filtered = [row for row in filtered if row.dossier_id.upper() == normalized]
    if changed_by:
        normalized_actor = changed_by.strip().lower()
        filtered = [row for row in filtered if row.changed_by.lower() == normalized_actor]
    if date_from:
        filtered = [row for row in filtered if row.changed_at.date() >= date_from]
    if date_to:
        filtered = [row for row in filtered if row.changed_at.date() <= date_to]
    return filtered


@router.get("/exports/indicators.csv")
async def export_indicators_csv(
    current_user: User = Depends(require_roles(Role.admin, Role.ministre)),
    db: Session = Depends(get_db),
):
    from ..main import _compute_sector_indicators

    indicators = _compute_sector_indicators(db)
    header = ["Secteur", "Volume local (T)", "Volume importe (T)", "Emplois"]
    rows = [
        [indicator.sector, indicator.local_volume_tons, indicator.import_volume_tons, indicator.jobs]
        for indicator in indicators
    ]
    write_audit_event(
        db,
        actor=current_user.username,
        action="export.csv",
        target="indicators",
        details="Export CSV indicateurs genere",
    )
    db.commit()
    return StreamingResponse(
        _csv_generator(rows, header),
        media_type="text/csv",
        headers={"Content-Disposition": 'attachment; filename="indicateurs-pnpi.csv"'},
    )


@router.get("/exports/dashboard.pdf")
async def export_dashboard_pdf(
    current_user: User = Depends(require_roles(Role.admin, Role.ministre)),
    db: Session = Depends(get_db),
) -> Response:
    from ..main import _compute_sector_indicators

    indicators = _compute_sector_indicators(db)
    all_units = db.execute(select(UnitORM)).scalars().unique().all()
    active_units = [unit for unit in all_units if unit.status == "active"]
    active_zones = len({unit.location for unit in active_units})
    traced_batches = db.execute(select(func.count(TraceBatchORM.batch_id))).scalar_one()
    total_local = sum(metric.local_volume_tons for metric in indicators)
    total_import = sum(metric.import_volume_tons for metric in indicators)
    denominator = total_local + total_import
    national_index = (total_local / denominator) if denominator > 0 else 0.0

    lines = [
        "REPUBLIQUE GABONAISE",
        "Ministere de l'Industrie et de la Transformation Locale",
        "CONFIDENTIEL · Document officiel PNPI",
        "",
        "PNPI - Resume Strategique",
        f"Date: {now_utc().isoformat()}",
        f"Indice national: {round(national_index * 100, 2)} %",
        f"Emplois industriels traces: {sum(metric.jobs for metric in indicators)}",
        f"Ecart import (T): {round(max(total_import - total_local, 0), 2)}",
        f"Unites actives: {len(active_units)}",
        f"Zones actives: {active_zones}",
        f"Lots traces: {traced_batches}",
    ]
    for indicator in indicators:
        lines.append(
            f"- {indicator.sector}: local {indicator.local_volume_tons}T / import {indicator.import_volume_tons}T / emplois {indicator.jobs}"
        )

    write_audit_event(
        db, actor=current_user.username, action="export.pdf", target="dashboard", details="Export PDF dashboard genere"
    )
    db.commit()

    return _build_pdf_response(lines, filename="pnpi-dashboard.pdf", font_size=10)


@router.get("/exports/inspectors-briefing.pdf")
async def export_inspectors_briefing_pdf(
    current_user: User = Depends(require_roles(Role.admin, Role.ministre)),
    db: Session = Depends(get_db),
) -> Response:
    field_reports = (
        db.execute(select(FieldReportORM).order_by(FieldReportORM.created_at.desc())).scalars().unique().all()
    )
    open_reports = [report for report in field_reports if report.status != "closed"]
    high_reports = [report for report in open_reports if report.severity.lower() in {"high", "critical"}]
    units = db.execute(select(UnitORM)).scalars().unique().all()

    lines = [
        "REPUBLIQUE GABONAISE",
        "Ministere de l'Industrie et de la Transformation Locale",
        "CONFIDENTIEL · Document officiel PNPI",
        "",
        "PNPI - Briefing Inspecteurs",
        f"Date: {now_utc().isoformat()}",
        f"Rapports terrain: {len(field_reports)}",
        f"Rapports ouverts: {len(open_reports)}",
        f"Rapports critiques ouverts: {len(high_reports)}",
        f"Unites suivies: {len(units)}",
    ]
    for report in high_reports[:8]:
        location = report.location or "Non renseigne"
        lines.append(f"- {report.id} [{report.status}] {report.severity} / {location} / {report.title}")

    write_audit_event(
        db,
        actor=current_user.username,
        action="export.pdf",
        target="inspectors-briefing",
        details="Export PDF briefing inspecteurs genere",
    )
    db.commit()

    return _build_pdf_response(lines, filename="pnpi-inspectors-briefing.pdf", font_size=10)


@router.get("/exports/pilotage-transitions.csv")
async def export_pilotage_transitions_csv(
    dossier_id: str | None = Query(default=None),
    changed_by: str | None = Query(default=None),
    date_from: date | None = Query(default=None),
    date_to: date | None = Query(default=None),
    current_user: User = Depends(require_roles(Role.admin, Role.ministre)),
    db: Session = Depends(get_db),
) -> Response:
    if date_from and date_to and date_from > date_to:
        raise HTTPException(status_code=400, detail="date_from doit etre <= date_to.")

    rows = (
        db.execute(select(ProjectDossierTransitionORM).order_by(ProjectDossierTransitionORM.changed_at.desc()))
        .scalars()
        .all()
    )
    rows = _filter_pilotage_transitions(
        rows,
        dossier_id=dossier_id,
        changed_by=changed_by,
        date_from=date_from,
        date_to=date_to,
    )
    header = [
        "Transition ID",
        "Dossier ID",
        "Acteur",
        "Date",
        "Statut precedent",
        "Statut nouveau",
        "Etape precedente",
        "Etape nouvelle",
        "Note",
    ]
    csv_rows = [
        [
            row.id,
            row.dossier_id,
            row.changed_by,
            row.changed_at.isoformat(),
            row.previous_status or "",
            row.new_status or "",
            row.previous_stage or "",
            row.new_stage or "",
            row.note,
        ]
        for row in rows
    ]
    write_audit_event(
        db,
        actor=current_user.username,
        action="exports.pilotage_transitions_csv",
        target=dossier_id.strip().upper() if dossier_id else None,
        details=(
            f"changed_by={changed_by or 'tous'}; date_from={date_from.isoformat() if date_from else 'Non renseigne'}; "
            f"date_to={date_to.isoformat() if date_to else 'Non renseigne'}; rows={len(rows)}"
        ),
    )
    db.commit()
    return StreamingResponse(
        _csv_generator(csv_rows, header),
        media_type="text/csv",
        headers={"Content-Disposition": 'attachment; filename="pilotage-transitions.csv"'},
    )


@router.get("/exports/pilotage-transitions.pdf")
async def export_pilotage_transitions_pdf(
    dossier_id: str | None = Query(default=None),
    changed_by: str | None = Query(default=None),
    date_from: date | None = Query(default=None),
    date_to: date | None = Query(default=None),
    current_user: User = Depends(require_roles(Role.admin, Role.ministre)),
    db: Session = Depends(get_db),
) -> Response:
    if date_from and date_to and date_from > date_to:
        raise HTTPException(status_code=400, detail="date_from doit etre <= date_to.")

    rows = (
        db.execute(select(ProjectDossierTransitionORM).order_by(ProjectDossierTransitionORM.changed_at.desc()))
        .scalars()
        .all()
    )
    rows = _filter_pilotage_transitions(
        rows,
        dossier_id=dossier_id,
        changed_by=changed_by,
        date_from=date_from,
        date_to=date_to,
    )

    lines = [
        "REPUBLIQUE GABONAISE",
        "Ministere de l'Industrie et de la Transformation Locale",
        "CONFIDENTIEL · Document officiel PNPI",
        "",
        "PNPI - Journal des transitions workflow",
        f"Date: {now_utc().isoformat()}",
        f"Nombre de transitions: {len(rows)}",
    ]
    for row in rows[:30]:
        lines.append(
            f"- {row.dossier_id} | {row.changed_by} | "
            f"{(row.previous_status or 'Non renseigne')}->{(row.new_status or 'Non renseigne')} | "
            f"{(row.previous_stage or 'Non renseigne')}->{(row.new_stage or 'Non renseigne')}"
        )
    write_audit_event(
        db,
        actor=current_user.username,
        action="exports.pilotage_transitions_pdf",
        target=dossier_id.strip().upper() if dossier_id else None,
        details=(
            f"changed_by={changed_by or 'tous'}; date_from={date_from.isoformat() if date_from else 'Non renseigne'}; "
            f"date_to={date_to.isoformat() if date_to else 'Non renseigne'}; rows={len(rows)}"
        ),
    )
    db.commit()

    return _build_pdf_response(
        lines, filename="pilotage-transitions.pdf", font_size=9, td_offset=12, start_y=770, start_x=36
    )


# ---------------------------------------------------------------------------
# PNPI CSV exports
# ---------------------------------------------------------------------------


@router.get("/pnpi/exports/ati.csv")
async def export_ati_csv(
    statut: str | None = Query(default=None),
    secteur: str | None = Query(default=None),
    current_user: User = Depends(
        require_roles(Role.admin, Role.ministre, Role.ministre, Role.directeur, Role.instructeur)
    ),
    db: Session = Depends(get_db),
):
    atis = (
        db.execute(
            select(AgrementTechniqueIndustrielORM)
            .where(*([] if not statut else [AgrementTechniqueIndustrielORM.statut == statut]))
            .where(*([] if not secteur else [AgrementTechniqueIndustrielORM.secteur == secteur]))
            .order_by(AgrementTechniqueIndustrielORM.date_soumission.desc())
        )
        .scalars()
        .all()
    )

    header = [
        "Numero ATI",
        "Operateur ID",
        "Type activite",
        "Secteur",
        "Statut",
        "Etape",
        "Priorite",
        "Instructeur",
        "Date soumission",
        "Date decision",
        "Age (j)",
        "SLA (j)",
        "En retard",
    ]
    rows = []
    for a in atis:
        age = max((now_utc().date() - a.date_soumission.date()).days, 0)
        overdue = age > a.sla_jours and a.statut not in {"approuve", "rejete", "expire"}
        rows.append(
            [
                a.numero_ati,
                a.operateur_id,
                a.type_activite,
                a.secteur,
                a.statut,
                a.etape,
                a.priorite,
                a.instructeur_username or "",
                a.date_soumission.date().isoformat(),
                a.date_decision.date().isoformat() if a.date_decision else "",
                age,
                a.sla_jours,
                "Oui" if overdue else "Non",
            ]
        )

    write_audit_event(
        db,
        actor=current_user.username,
        action="export.csv",
        target="ati",
        details=f"Export CSV ATI genere ({len(rows)} lignes)",
    )
    db.commit()

    return StreamingResponse(
        _csv_generator(rows, header),
        media_type="text/csv",
        headers={"Content-Disposition": 'attachment; filename="ati_export.csv"'},
    )


@router.get("/pnpi/exports/operateurs.csv")
async def export_operateurs_csv(
    secteur: str | None = Query(default=None),
    current_user: User = Depends(
        require_roles(Role.admin, Role.ministre, Role.ministre, Role.directeur, Role.instructeur)
    ),
    db: Session = Depends(get_db),
):
    ops = (
        db.execute(
            select(OperateurIndustrielORM)
            .where(*([] if not secteur else [OperateurIndustrielORM.secteur == secteur]))
            .order_by(OperateurIndustrielORM.raison_sociale)
        )
        .scalars()
        .all()
    )

    header = [
        "ID",
        "NIF Gabon",
        "Raison sociale",
        "Secteur",
        "Province",
        "Ville",
        "Effectif declare",
        "Statut",
        "Email",
        "Telephone",
        "Cree le",
    ]
    rows = []
    for op in ops:
        rows.append(
            [
                op.id,
                op.nif,
                op.raison_sociale,
                op.secteur,
                op.province,
                op.ville,
                op.effectif_declare or 0,
                "Actif" if op.is_active else "Inactif",
                op.contact_email or "",
                op.contact_telephone or "",
                op.created_at.date().isoformat(),
            ]
        )

    write_audit_event(
        db,
        actor=current_user.username,
        action="export.csv",
        target="operateurs",
        details=f"Export CSV operateurs genere ({len(rows)} lignes)",
    )
    db.commit()

    return StreamingResponse(
        _csv_generator(rows, header),
        media_type="text/csv",
        headers={"Content-Disposition": 'attachment; filename="operateurs_export.csv"'},
    )


@router.get("/pnpi/exports/inspections.csv", summary="Export CSV des inspections de conformite")
async def export_inspections_csv(
    statut_conformite: str | None = Query(default=None),
    current_user: User = Depends(
        require_roles(Role.admin, Role.ministre, Role.directeur, Role.instructeur, Role.inspecteur)
    ),
    db: Session = Depends(get_db),
):
    query = select(InspectionConformiteORM).order_by(InspectionConformiteORM.date_inspection.desc())
    if statut_conformite:
        query = query.where(InspectionConformiteORM.statut_conformite == statut_conformite)
    inspections = db.execute(query).scalars().all()

    header = [
        "ID",
        "Operateur ID",
        "ATI ID",
        "Inspecteur",
        "Date",
        "Statut conformite",
        "Observations",
        "Mesures correctives",
        "Latitude",
        "Longitude",
    ]
    rows = []
    for insp in inspections:
        rows.append(
            [
                insp.id,
                insp.operateur_id,
                insp.ati_id or "",
                insp.inspecteur_username,
                insp.date_inspection.date().isoformat(),
                insp.statut_conformite,
                insp.observations[:200] if insp.observations else "",
                (insp.mesures_correctives or "")[:200],
                insp.latitude or "",
                insp.longitude or "",
            ]
        )

    write_audit_event(
        db,
        actor=current_user.username,
        action="export.csv",
        target="inspections",
        details=f"Export CSV inspections genere ({len(rows)} lignes)",
    )
    db.commit()

    return StreamingResponse(
        _csv_generator(rows, header),
        media_type="text/csv",
        headers={"Content-Disposition": 'attachment; filename="inspections_export.csv"'},
    )


@router.get("/pnpi/exports/briefing.pdf", summary="Briefing ministeriel PNPI (PDF)")
async def export_pnpi_briefing_pdf(
    current_user: User = Depends(require_roles(Role.admin, Role.ministre, Role.directeur)),
    db: Session = Depends(get_db),
) -> Response:
    """Genere un PDF de synthese ministerielle : KPIs, pipeline, repartition sectorielle."""
    from collections import defaultdict as _defaultdict
    from statistics import median as _median

    all_atis = db.execute(select(AgrementTechniqueIndustrielORM)).scalars().all()
    all_ops = db.execute(select(OperateurIndustrielORM)).scalars().all()

    _TERM = {"approuve", "rejete", "expire"}
    now = now_utc()
    first_of_month = now.date().replace(day=1)

    # KPIs
    atis_total = len(all_atis)
    atis_en_cours = sum(1 for a in all_atis if a.statut not in _TERM)
    atis_approuves_mois = sum(
        1 for a in all_atis if a.statut == "approuve" and a.date_decision and a.date_decision.date() >= first_of_month
    )
    atis_en_retard = sum(
        1 for a in all_atis if a.statut not in _TERM and (now.date() - a.date_soumission.date()).days > a.sla_jours
    )
    decided = [a for a in all_atis if a.statut in {"approuve", "rejete"} and a.date_decision]
    durations = [max((a.date_decision.date() - a.date_soumission.date()).days, 0) for a in decided]
    delai_moyen = float(_median(durations)) if durations else 0.0
    compliant = sum(1 for a in decided if (a.date_decision.date() - a.date_soumission.date()).days <= a.sla_jours)
    taux_sla = round((compliant / len(durations) * 100) if durations else 0.0, 0)
    ops_actifs = sum(1 for op in all_ops if op.is_active)

    # Pipeline
    counts: dict[str, int] = _defaultdict(int)
    for a in all_atis:
        counts[a.statut] += 1

    # Secteurs
    sec_ops: dict[str, int] = _defaultdict(int)
    sec_atis: dict[str, int] = _defaultdict(int)
    sec_approuves: dict[str, int] = _defaultdict(int)
    sec_emplois: dict[str, int] = _defaultdict(int)
    for op in all_ops:
        sec_ops[op.secteur] += 1
        sec_emplois[op.secteur] += op.effectif_declare or 0
    for a in all_atis:
        sec_atis[a.secteur] += 1
        if a.statut == "approuve":
            sec_approuves[a.secteur] += 1
    all_secteurs = sorted(set(list(sec_ops.keys()) + list(sec_atis.keys())))

    # Provinces
    prov_ops: dict[str, int] = _defaultdict(int)
    prov_atis: dict[str, int] = _defaultdict(int)
    op_prov: dict[str, str] = {}
    for op in all_ops:
        prov_ops[op.province] += 1
        op_prov[op.id] = op.province
    for a in all_atis:
        if a.statut not in _TERM:
            prov_atis[op_prov.get(a.operateur_id, "inconnu")] += 1
    all_provinces = sorted(set(list(prov_ops.keys()) + list(prov_atis.keys())))

    ts = now_utc().strftime("%d/%m/%Y %H:%M")
    lines = [
        "REPUBLIQUE GABONAISE",
        "Ministere de l'Industrie et de la Transformation Locale",
        "CONFIDENTIEL · Document officiel PNPI",
        "",
        f"BRIEFING MINISTERIEL PNPI · {ts}",
        "=" * 60,
        "",
        "INDICATEURS CLES",
        f"  ATIs total .............. {atis_total}",
        f"  ATIs en cours ........... {atis_en_cours}",
        f"  ATIs approuves ce mois .. {atis_approuves_mois}",
        f"  ATIs en retard SLA ...... {atis_en_retard}",
        f"  Delai moyen (jours) ..... {delai_moyen:.1f}",
        f"  Taux conformite SLA ..... {taux_sla:.0f} %",
        f"  Operateurs actifs ....... {ops_actifs}",
        "",
        "PIPELINE DES STATUTS",
        f"  Soumis .................. {counts.get('soumis', 0)}",
        f"  En instruction .......... {counts.get('en_instruction', 0)}",
        f"  En validation ........... {counts.get('en_validation', 0)}",
        f"  Approuves ............... {counts.get('approuve', 0)}",
        f"  Rejetes ................. {counts.get('rejete', 0)}",
        f"  Expires ................. {counts.get('expire', 0)}",
        "",
        "REPARTITION PAR SECTEUR",
    ]
    for s in all_secteurs:
        nb = sec_atis.get(s, 0)
        taux = f"{round(sec_approuves.get(s, 0) / nb * 100):.0f}%" if nb > 0 else "N/A"
        lines.append(
            f"  {s:<18} {sec_ops.get(s, 0):>3} ops  {nb:>3} ATIs  {sec_emplois.get(s, 0):>5} emplois  taux {taux}"
        )
    lines.append("")
    lines.append("REPARTITION PAR PROVINCE")
    for p in all_provinces:
        lines.append(
            f"  {p.replace('_', ' ').title():<22} {prov_ops.get(p, 0):>3} ops  {prov_atis.get(p, 0):>3} ATIs actifs"
        )
    lines.append("")
    lines.append(f"Document genere le {ts} par {current_user.username}")

    write_audit_event(
        db,
        actor=current_user.username,
        action="exports.pnpi_briefing_pdf",
        target=None,
        details=f"atis_total={atis_total}; approuves={counts.get('approuve', 0)}",
    )
    db.commit()

    return _build_pdf_response(
        lines, filename="pnpi-briefing-ministeriel.pdf", font_size=9, td_offset=12, start_y=770, start_x=36
    )


@router.post("/admin/briefing/email")
async def email_briefing_to_roles(
    target_roles: list[str] = Body(default=["ministre", "directeur"]),
    current_user: User = Depends(require_roles(Role.admin, Role.ministre)),
    db: Session = Depends(get_db),
):
    """Envoie un email de briefing aux utilisateurs ayant les roles cibles."""
    from ..core.email import send_email

    # Find users with target roles
    all_users = db.execute(select(UserAccountORM).where(UserAccountORM.is_active.is_(True))).scalars().all()
    recipients: list[str] = []
    for user in all_users:
        user_roles = set(user.roles_csv.split(","))
        if user_roles & set(target_roles):
            # Use username as email placeholder (format: username@pnpi-gabon.ga)
            recipients.append(f"{user.username}@pnpi-gabon.ga")

    # Generate briefing content
    now = now_utc()
    subject = f"Briefing PNPI \u2014 {now.strftime('%d/%m/%Y')}"
    html = f"""
    <div style="font-family: Arial, sans-serif; max-width: 600px; margin: 0 auto;">
      <div style="background: #0d3f78; color: white; padding: 20px; border-radius: 12px 12px 0 0;">
        <h1 style="margin: 0; font-size: 20px;">PNPI \u2014 Briefing Quotidien</h1>
      </div>
      <div style="padding: 24px; background: #f6f8fb; border-radius: 0 0 12px 12px;">
        <p>Le briefing du {now.strftime("%d/%m/%Y")} est disponible sur la plateforme.</p>
        <a href="https://pnpi-gabon.ga/briefing" style="display: inline-block; padding: 12px 24px; background: #006233; color: white; border-radius: 8px; text-decoration: none; font-weight: bold;">
          Consulter le briefing
        </a>
      </div>
    </div>
    """

    sent = 0
    if recipients:
        for email_addr in recipients:
            if send_email([email_addr], subject, html):
                sent += 1

    write_audit_event(
        db,
        actor=current_user.username,
        action="briefing.email",
        target=f"{len(target_roles)} roles",
        details=f"Briefing envoye a {sent} destinataire(s)",
    )
    db.commit()

    return {"sent": sent, "recipients_found": len(recipients), "target_roles": target_roles}


@router.get("/exports/ati/{ati_id}/documents.zip")
async def export_ati_documents_zip(
    ati_id: str,
    current_user: User = Depends(
        require_roles(Role.admin, Role.ministre, Role.directeur, Role.instructeur, Role.operateur)
    ),
    db: Session = Depends(get_db),
):
    """Download all documents for an ATI as a ZIP archive."""
    from .ati import check_ati_access

    ati = db.get(AgrementTechniqueIndustrielORM, ati_id)
    if not ati:
        raise HTTPException(status_code=404, detail="ATI introuvable.")
    check_ati_access(ati, current_user)

    docs = db.execute(select(DocumentDossierORM).where(DocumentDossierORM.ati_id == ati_id)).scalars().all()

    if not docs:
        raise HTTPException(status_code=404, detail="Aucun document pour cet ATI.")

    storage = get_storage(os.getenv("PNPI_UPLOAD_DIR", "uploads/ati"))
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for doc in docs:
            if storage.exists(doc.chemin_stockage):
                zf.writestr(doc.nom_fichier, storage.read(doc.chemin_stockage))

    buf.seek(0)

    write_audit_event(
        db,
        actor=current_user.username,
        action="export.zip",
        target=ati_id,
        details=f"ZIP {len(docs)} documents ATI {ati.numero_ati}",
    )
    db.commit()

    return Response(
        content=buf.read(),
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="documents_ATI_{ati.numero_ati}.zip"'},
    )


# ---------------------------------------------------------------------------
# Excel (XLSX) exports
# ---------------------------------------------------------------------------


@router.get("/exports/atis.xlsx")
async def export_atis_xlsx(
    statut: str | None = Query(None),
    secteur: str | None = Query(None),
    province: str | None = Query(None),
    current_user: User = Depends(require_roles(Role.admin, Role.ministre, Role.directeur, Role.instructeur)),
    db: Session = Depends(get_db),
):
    """Export ATI list as Excel workbook with formatting."""
    query = select(AgrementTechniqueIndustrielORM)
    if statut:
        query = query.where(AgrementTechniqueIndustrielORM.statut == statut)
    if secteur:
        query = query.where(AgrementTechniqueIndustrielORM.secteur == secteur)

    atis = db.execute(query.order_by(AgrementTechniqueIndustrielORM.date_soumission.desc())).scalars().all()

    if province:
        atis = [a for a in atis if a.operateur and a.operateur.province == province]

    wb = Workbook()
    ws = wb.active
    ws.title = "ATIs PNPI"

    # Header style
    header_font = Font(name="Calibri", bold=True, size=11, color="FFFFFF")
    header_fill = PatternFill(start_color="006233", end_color="006233", fill_type="solid")
    header_align = Alignment(horizontal="center", vertical="center")
    thin_border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin"),
    )

    headers = [
        "N\u00b0 ATI",
        "Operateur",
        "NIF",
        "Secteur",
        "Province",
        "Type activite",
        "Statut",
        "Date soumission",
        "Date decision",
        "SLA (jours)",
        "Age (jours)",
    ]

    for col, header in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = header_align
        cell.border = thin_border

    # Status colors
    STATUS_FILLS = {
        "approuve": PatternFill(start_color="DCFCE7", end_color="DCFCE7", fill_type="solid"),
        "rejete": PatternFill(start_color="FEF2F2", end_color="FEF2F2", fill_type="solid"),
        "soumis": PatternFill(start_color="F3F4F6", end_color="F3F4F6", fill_type="solid"),
    }

    from datetime import datetime

    now = datetime.now(UTC)

    for row_idx, ati in enumerate(atis, 2):
        op = ati.operateur
        age = (now.date() - ati.date_soumission.date()).days

        values = [
            ati.numero_ati,
            op.raison_sociale if op else "",
            (op.nif or "") if op else "",
            ati.secteur,
            (op.province if op else "").replace("_", " ").title(),
            ati.type_activite,
            ati.statut,
            ati.date_soumission.strftime("%d/%m/%Y"),
            ati.date_decision.strftime("%d/%m/%Y") if ati.date_decision else "",
            ati.sla_jours,
            age,
        ]

        for col, val in enumerate(values, 1):
            cell = ws.cell(row=row_idx, column=col, value=val)
            cell.border = thin_border
            cell.alignment = Alignment(vertical="center")

        # Color status cell
        status_cell = ws.cell(row=row_idx, column=7)
        if ati.statut in STATUS_FILLS:
            status_cell.fill = STATUS_FILLS[ati.statut]

        # Highlight overdue
        if age > ati.sla_jours and ati.statut not in ("approuve", "rejete", "expire"):
            ws.cell(row=row_idx, column=11).fill = PatternFill(
                start_color="FEF2F2", end_color="FEF2F2", fill_type="solid"
            )
            ws.cell(row=row_idx, column=11).font = Font(color="B42318", bold=True)

    # Auto-width columns
    for col in ws.columns:
        max_length = max(len(str(cell.value or "")) for cell in col)
        ws.column_dimensions[col[0].column_letter].width = min(max_length + 4, 40)

    # Freeze header
    ws.freeze_panes = "A2"

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)

    write_audit_event(
        db, actor=current_user.username, action="export.xlsx", target="atis", details=f"Export Excel {len(atis)} ATIs"
    )
    db.commit()

    return Response(
        content=buf.read(),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": 'attachment; filename="atis_pnpi.xlsx"'},
    )


@router.get("/exports/operateurs.xlsx")
async def export_operateurs_xlsx(
    secteur: str | None = Query(None),
    province: str | None = Query(None),
    current_user: User = Depends(require_roles(Role.admin, Role.ministre, Role.directeur)),
    db: Session = Depends(get_db),
):
    """Export operators as Excel workbook."""
    query = select(OperateurIndustrielORM).where(OperateurIndustrielORM.is_active.is_(True))
    if secteur:
        query = query.where(OperateurIndustrielORM.secteur == secteur)
    if province:
        query = query.where(OperateurIndustrielORM.province == province)

    ops = db.execute(query.order_by(OperateurIndustrielORM.raison_sociale)).scalars().all()

    wb = Workbook()
    ws = wb.active
    ws.title = "Operateurs PNPI"

    header_font = Font(name="Calibri", bold=True, size=11, color="FFFFFF")
    header_fill = PatternFill(start_color="0C7EB4", end_color="0C7EB4", fill_type="solid")
    thin_border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin"),
    )

    headers = ["Raison sociale", "NIF", "Secteur", "Province", "Ville", "Email", "Telephone"]
    for col, header in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = thin_border

    for row_idx, op in enumerate(ops, 2):
        values = [
            op.raison_sociale,
            op.nif,
            op.secteur,
            op.province.replace("_", " ").title() if op.province else "",
            op.ville or "",
            op.email or "",
            op.telephone or "",
        ]
        for col, val in enumerate(values, 1):
            cell = ws.cell(row=row_idx, column=col, value=val)
            cell.border = thin_border

    for col in ws.columns:
        max_length = max(len(str(cell.value or "")) for cell in col)
        ws.column_dimensions[col[0].column_letter].width = min(max_length + 4, 40)

    ws.freeze_panes = "A2"

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)

    write_audit_event(
        db,
        actor=current_user.username,
        action="export.xlsx",
        target="operateurs",
        details=f"Export Excel {len(ops)} operateurs",
    )
    db.commit()

    return Response(
        content=buf.read(),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": 'attachment; filename="operateurs_pnpi.xlsx"'},
    )


# ---------------------------------------------------------------------------
# PowerPoint (PPTX) export
# ---------------------------------------------------------------------------


@router.get("/exports/briefing.pptx")
async def export_briefing_pptx(
    current_user: User = Depends(require_roles(Role.admin, Role.ministre, Role.directeur)),
    db: Session = Depends(get_db),
):
    """Export executive briefing as PowerPoint presentation."""

    # Fetch data
    all_atis = db.execute(select(AgrementTechniqueIndustrielORM)).scalars().all()
    all_ops = (
        db.execute(select(OperateurIndustrielORM).where(OperateurIndustrielORM.is_active.is_(True))).scalars().all()
    )
    all_insp = db.execute(select(InspectionConformiteORM)).scalars().all()

    now = now_utc()

    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    VERT = RGBColor(0x00, 0x62, 0x33)
    BLEU = RGBColor(0x05, 0x1B, 0x36)
    BLANC = RGBColor(0xFF, 0xFF, 0xFF)
    GRIS = RGBColor(0x52, 0x61, 0x75)

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BLEU
        # Title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = BLANC
        p.alignment = PP_ALIGN.CENTER
        # Subtitle
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0x9E, 0xB4, 0xCF)
        p2.alignment = PP_ALIGN.CENTER

    def add_kpi_slide(title, kpis_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BLANC
        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BLEU
        # KPI boxes
        cols = min(len(kpis_data), 4)
        box_w = 2.8
        gap = 0.3
        start_x = (13.333 - (cols * box_w + (cols - 1) * gap)) / 2
        for i, kpi in enumerate(kpis_data[:4]):
            x = start_x + i * (box_w + gap)
            shape = slide.shapes.add_textbox(Inches(x), Inches(1.8), Inches(box_w), Inches(2.2))
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(kpi["value"])
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = VERT
            p.alignment = PP_ALIGN.CENTER
            p2 = tf.add_paragraph()
            p2.text = kpi["label"]
            p2.font.size = Pt(14)
            p2.font.color.rgb = GRIS
            p2.alignment = PP_ALIGN.CENTER

    # Slide 1: Title
    add_title_slide(
        "PNPI \u2014 Briefing Executif",
        f"Ministere de l'Industrie et de la Transformation Locale\n{now.strftime('%d/%m/%Y')}",
    )

    # Slide 2: ATI KPIs
    approved = sum(1 for a in all_atis if a.statut == "approuve")
    rejected = sum(1 for a in all_atis if a.statut == "rejete")
    in_progress = sum(1 for a in all_atis if a.statut not in ("approuve", "rejete", "expire"))
    add_kpi_slide(
        "Agrements Techniques Industriels",
        [
            {"label": "Total ATIs", "value": len(all_atis)},
            {"label": "Approuves", "value": approved},
            {"label": "Rejetes", "value": rejected},
            {"label": "En cours", "value": in_progress},
        ],
    )

    # Slide 3: Operators & Inspections
    conformes = sum(1 for i in all_insp if i.statut_conformite == "conforme")
    add_kpi_slide(
        "Operateurs & Inspections",
        [
            {"label": "Operateurs actifs", "value": len(all_ops)},
            {"label": "Inspections", "value": len(all_insp)},
            {"label": "Conformes", "value": conformes},
            {"label": "Taux conformite", "value": f"{round(conformes / len(all_insp) * 100)}%" if all_insp else "0%"},
        ],
    )

    # Slide 4: SLA
    terminal = {"approuve", "rejete", "expire"}
    overdue = sum(
        1 for a in all_atis if a.statut not in terminal and (now.date() - a.date_soumission.date()).days > a.sla_jours
    )
    decided = [a for a in all_atis if a.statut in ("approuve", "rejete") and a.date_decision]
    avg_days = (
        round(sum((a.date_decision.date() - a.date_soumission.date()).days for a in decided) / len(decided))
        if decided
        else 0
    )
    add_kpi_slide(
        "Performance SLA",
        [
            {"label": "Depassements SLA", "value": overdue},
            {"label": "Delai moyen (jours)", "value": avg_days},
            {
                "label": "Taux respect SLA",
                "value": f"{round((1 - overdue / len(all_atis)) * 100)}%" if all_atis else "100%",
            },
        ],
    )

    # Slide 5: Closing
    add_title_slide("Merci", "Questions ?\npnpi-gabon.ga")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    write_audit_event(
        db,
        actor=current_user.username,
        action="export.pptx",
        target="briefing",
        details="Export PowerPoint briefing executif",
    )
    db.commit()

    return Response(
        content=buf.read(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="briefing_pnpi_{now.strftime("%Y%m%d")}.pptx"'},
    )


# ---------------------------------------------------------------------------
# Signed certificate & signature verification
# ---------------------------------------------------------------------------


@router.get("/exports/ati/{ati_id}/signed-certificate.pdf")
async def export_signed_certificate(
    ati_id: str,
    current_user: User = Depends(require_roles(Role.admin, Role.ministre, Role.directeur)),
    db: Session = Depends(get_db),
):
    """Generate a digitally signed ATI certificate."""
    ati = db.get(AgrementTechniqueIndustrielORM, ati_id)
    if not ati or ati.statut != "approuve":
        raise HTTPException(400, "Certificat disponible uniquement pour les ATI approuves.")

    # Generate certificate PDF
    from ..core.certificate import generate_ati_certificate

    op = ati.operateur
    pdf_content = generate_ati_certificate(
        numero_ati=ati.numero_ati,
        operateur=op.raison_sociale if op else "Inconnu",
        nif=(op.nif or "") if op else "",
        secteur=ati.secteur,
        province=op.province if op else "",
        type_activite=ati.type_activite,
        date_soumission=ati.date_soumission,
        date_decision=ati.date_decision,
        date_expiration=ati.date_expiration,
        reference_decision=ati.numero_reference_decision if hasattr(ati, "numero_reference_decision") else None,
        sla_jours=ati.sla_jours,
    )

    # Generate digital signature
    sig = generate_signature(
        document_content=pdf_content,
        signer_username=current_user.username,
        signer_role=current_user.roles[0] if current_user.roles else "admin",
        document_type="certificat_ati",
        reference=ati.numero_ati,
    )

    # Store signature in audit log
    write_audit_event(
        db,
        actor=current_user.username,
        action="ati.sign_certificate",
        target=ati.id,
        details=f"Certificat signe: {sig['signature'][:16]}... par {current_user.username}",
    )
    db.commit()

    return Response(
        content=pdf_content,
        media_type="application/pdf",
        headers={
            "Content-Disposition": f'attachment; filename="certificat_signe_ATI_{ati.numero_ati}.pdf"',
            "X-PNPI-Signature": sig["signature"],
            "X-PNPI-Signer": sig["signer"],
            "X-PNPI-Signed-At": sig["timestamp"],
        },
    )


@router.post("/exports/verify-signature")
async def verify_document_signature(
    data: dict,
):
    """Public endpoint to verify a document signature."""
    from ..core.signature import verify_signature

    # In production, would accept file upload + signature data
    signature_data = data.get("signature_data", {})
    content_hash = data.get("content_hash", "")

    if not signature_data:
        raise HTTPException(400, "Donnees de signature requises.")

    # Simple hash-based verification (no file upload needed)
    if content_hash:
        signature_data["content_hash"] = content_hash

    result = verify_signature(signature_data, bytes.fromhex("00"))  # Placeholder
    return result


@router.get("/exports/batch-qr.pdf")
async def export_batch_qr_pdf(
    statut: str = Query("approuve"),
    current_user: User = Depends(require_roles(Role.admin, Role.ministre, Role.directeur)),
    db: Session = Depends(get_db),
):
    """Generate a PDF with QR codes for all approved ATIs (or filtered by status)."""
    query = (
        select(AgrementTechniqueIndustrielORM)
        .where(AgrementTechniqueIndustrielORM.statut == statut)
        .order_by(AgrementTechniqueIndustrielORM.numero_ati)
    )

    atis = db.execute(query).scalars().all()

    if not atis:
        raise HTTPException(404, f"Aucun ATI avec le statut '{statut}'.")

    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import Image as RLImage
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4, leftMargin=1.5 * cm, rightMargin=1.5 * cm, topMargin=2 * cm, bottomMargin=1.5 * cm
    )
    styles = getSampleStyleSheet()
    center = ParagraphStyle("center", parent=styles["Normal"], alignment=TA_CENTER, fontSize=10)

    story = []
    story.append(
        Paragraph(
            "PNPI · QR Codes de Verification",
            ParagraphStyle(
                "title", parent=center, fontSize=16, fontName="Helvetica-Bold", textColor=colors.HexColor("#006233")
            ),
        )
    )
    story.append(
        Paragraph(
            f"Statut: {statut} · {len(atis)} ATI(s)",
            ParagraphStyle("sub", parent=center, fontSize=10, textColor=colors.gray),
        )
    )
    story.append(Spacer(1, 0.5 * cm))

    # Build rows of 3 QR codes
    rows = []
    current_row = []

    for ati in atis:
        qr_url = f"https://pnpi-gabon.ga/verify/ati/{ati.numero_ati}"
        qr = qrcode.QRCode(version=1, box_size=6, border=2)
        qr.add_data(qr_url)
        qr.make(fit=True)
        qr_img = qr.make_image(fill_color="#003F8F", back_color="white")
        qr_buf = io.BytesIO()
        qr_img.save(qr_buf, format="PNG")
        qr_buf.seek(0)

        op_name = ati.operateur.raison_sociale[:25] if ati.operateur else "·"

        cell_content = [
            RLImage(qr_buf, width=3.5 * cm, height=3.5 * cm),
            Paragraph(
                f"<b>{ati.numero_ati}</b>",
                ParagraphStyle("qr_num", parent=center, fontSize=8, fontName="Helvetica-Bold"),
            ),
            Paragraph(op_name, ParagraphStyle("qr_op", parent=center, fontSize=7, textColor=colors.gray)),
        ]

        current_row.append(cell_content)

        if len(current_row) == 3:
            rows.append(current_row)
            current_row = []

    if current_row:
        while len(current_row) < 3:
            current_row.append(["", "", ""])
        rows.append(current_row)

    for row in rows:
        t = Table(
            [
                [c[0] if c else "" for c in row],
                [c[1] if c and len(c) > 1 else "" for c in row],
                [c[2] if c and len(c) > 2 else "" for c in row],
            ],
            colWidths=[6 * cm, 6 * cm, 6 * cm],
        )
        t.setStyle(
            TableStyle(
                [
                    ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ("TOPPADDING", (0, 0), (-1, -1), 8),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ]
            )
        )
        story.append(t)
        story.append(Spacer(1, 0.3 * cm))

    doc.build(story)
    buf.seek(0)

    write_audit_event(
        db,
        actor=current_user.username,
        action="export.batch_qr",
        target=statut,
        details=f"Batch QR PDF {len(atis)} ATIs",
    )
    db.commit()

    return Response(
        content=buf.read(),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="qr_codes_ati_{statut}.pdf"'},
    )


# ---------------------------------------------------------------------------
# PDF builder helper
# ---------------------------------------------------------------------------


def _build_pdf_response(
    lines: list,
    *,
    filename: str,
    font_size: int = 10,
    td_offset: int = 14,
    start_y: int = 760,
    start_x: int = 40,
) -> Response:
    escaped_lines = [line.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)") for line in lines]
    text_commands = [f"BT /F1 {font_size} Tf {start_x} {start_y} Td"]
    for index, line in enumerate(escaped_lines):
        if index > 0:
            text_commands.append(f"0 -{td_offset} Td")
        text_commands.append(f"({line}) Tj")
    text_commands.append("ET")
    content_stream = "\n".join(text_commands).encode("latin-1", errors="replace")

    objects_list = [
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n",
        b"2 0 obj\n<< /Type /Pages /Count 1 /Kids [3 0 R] >>\nendobj\n",
        b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\nendobj\n",
        f"4 0 obj\n<< /Length {len(content_stream)} >>\nstream\n".encode("latin-1")
        + content_stream
        + b"\nendstream\nendobj\n",
        b"5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n",
    ]

    pdf_body = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for obj in objects_list:
        offsets.append(len(pdf_body))
        pdf_body.extend(obj)
    xref_offset = len(pdf_body)
    pdf_body.extend(f"xref\n0 {len(offsets)}\n".encode("latin-1"))
    pdf_body.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        pdf_body.extend(f"{offset:010d} 00000 n \n".encode("latin-1"))
    pdf_body.extend(
        f"trailer\n<< /Size {len(offsets)} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF".encode("latin-1")
    )
    return Response(
        content=bytes(pdf_body),
        media_type="application/pdf",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


@router.get("/exports/ati/{ati_id}/letter.pdf")
async def generate_official_letter(
    ati_id: str,
    letter_type: str = Query("approval", description="approval|rejection|acknowledgment"),
    current_user: User = Depends(require_roles(Role.admin, Role.ministre, Role.directeur)),
    db: Session = Depends(get_db),
):
    """Generate an official ministerial letter for an ATI."""
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import HRFlowable, Paragraph, SimpleDocTemplate, Spacer

    ati = db.get(AgrementTechniqueIndustrielORM, ati_id)
    if not ati:
        raise HTTPException(404, "ATI introuvable.")
    op = ati.operateur

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4, leftMargin=3 * cm, rightMargin=3 * cm, topMargin=2.5 * cm, bottomMargin=2.5 * cm
    )
    styles = getSampleStyleSheet()

    center = ParagraphStyle("center", parent=styles["Normal"], alignment=TA_CENTER, fontSize=10)
    justify = ParagraphStyle("justify", parent=styles["Normal"], alignment=TA_JUSTIFY, fontSize=11, leading=16)

    story = []

    # Header
    story.append(
        Paragraph("REPUBLIQUE GABONAISE", ParagraphStyle("rg", parent=center, fontSize=11, textColor=colors.gray))
    )
    story.append(
        Paragraph(
            "Union \u2014 Travail \u2014 Justice",
            ParagraphStyle("motto", parent=center, fontSize=9, textColor=colors.gray, fontName="Helvetica-Oblique"),
        )
    )
    story.append(Spacer(1, 0.3 * cm))
    story.append(
        Paragraph(
            "MINISTERE DE L'INDUSTRIE ET DE LA TRANSFORMATION LOCALE",
            ParagraphStyle(
                "min", parent=center, fontSize=9, textColor=colors.HexColor("#003F8F"), fontName="Helvetica-Bold"
            ),
        )
    )
    story.append(
        Paragraph(
            "Direction Generale de l'Industrie",
            ParagraphStyle("dgi", parent=center, fontSize=8, textColor=colors.HexColor("#003F8F")),
        )
    )
    story.append(Spacer(1, 0.3 * cm))
    story.append(HRFlowable(width="100%", thickness=2, color=colors.HexColor("#006233")))
    story.append(Spacer(1, 0.8 * cm))

    now = now_utc()
    story.append(
        Paragraph(
            f"Libreville, le {now.strftime('%d/%m/%Y')}",
            ParagraphStyle("date", parent=styles["Normal"], fontSize=10, alignment=2),
        )
    )
    story.append(Spacer(1, 0.5 * cm))

    # Recipient
    op_name = op.raison_sociale if op else "Operateur"
    story.append(
        Paragraph(f"<b>A l'attention de :</b> {op_name}", ParagraphStyle("to", parent=styles["Normal"], fontSize=11))
    )
    story.append(Spacer(1, 0.3 * cm))
    story.append(
        Paragraph(
            f"<b>Objet :</b> Dossier ATI N\u00b0 {ati.numero_ati}",
            ParagraphStyle("obj", parent=styles["Normal"], fontSize=11),
        )
    )
    story.append(Spacer(1, 0.5 * cm))

    # Body
    if letter_type == "approval":
        story.append(Paragraph("Madame, Monsieur,", justify))
        story.append(Spacer(1, 0.3 * cm))
        story.append(
            Paragraph(
                f"J'ai l'honneur de vous informer que votre demande d'Agrement Technique Industriel "
                f"N\u00b0 <b>{ati.numero_ati}</b>, deposee le {ati.date_soumission.strftime('%d/%m/%Y')}, "
                f"portant sur l'activite de <i>{ati.type_activite[:100]}</i> dans le secteur <b>{ati.secteur}</b>, "
                f"a ete examinee favorablement par nos services.",
                justify,
            )
        )
        story.append(Spacer(1, 0.3 * cm))
        story.append(
            Paragraph(
                "En consequence, votre agrement vous est accorde conformement aux dispositions "
                "reglementaires en vigueur. Vous trouverez ci-joint votre certificat d'approbation.",
                justify,
            )
        )
    elif letter_type == "rejection":
        story.append(Paragraph("Madame, Monsieur,", justify))
        story.append(Spacer(1, 0.3 * cm))
        story.append(
            Paragraph(
                f"J'ai le regret de vous informer que votre demande d'Agrement Technique Industriel "
                f"N\u00b0 <b>{ati.numero_ati}</b>, deposee le {ati.date_soumission.strftime('%d/%m/%Y')}, "
                f"n'a pas pu recevoir une suite favorable en l'etat actuel de votre dossier.",
                justify,
            )
        )
        story.append(Spacer(1, 0.3 * cm))
        story.append(
            Paragraph(
                "Nous vous invitons a prendre connaissance des observations formulees par nos services "
                "et a resoumettre votre dossier apres avoir apporte les corrections necessaires.",
                justify,
            )
        )
    else:
        story.append(Paragraph("Madame, Monsieur,", justify))
        story.append(Spacer(1, 0.3 * cm))
        story.append(
            Paragraph(
                f"J'accuse reception de votre demande d'Agrement Technique Industriel "
                f"N\u00b0 <b>{ati.numero_ati}</b>, deposee le {ati.date_soumission.strftime('%d/%m/%Y')}. "
                f"Votre dossier est en cours d'instruction par nos services.",
                justify,
            )
        )

    story.append(Spacer(1, 0.3 * cm))
    story.append(Paragraph("Veuillez agreer, Madame, Monsieur, l'expression de ma consideration distinguee.", justify))

    story.append(Spacer(1, 1.5 * cm))
    story.append(
        Paragraph(
            "<b>Le Directeur General de l'Industrie</b>",
            ParagraphStyle("sig", parent=styles["Normal"], fontSize=10, alignment=2),
        )
    )

    story.append(Spacer(1, 2 * cm))
    story.append(HRFlowable(width="100%", thickness=0.5, color=colors.gray))
    story.append(
        Paragraph(
            "Document genere par la Plateforme Nationale de la Politique Industrielle (PNPI)",
            ParagraphStyle("footer", parent=center, fontSize=7, textColor=colors.gray),
        )
    )

    doc.build(story)
    buf.seek(0)

    write_audit_event(
        db,
        actor=current_user.username,
        action="export.letter",
        target=ati_id,
        details=f"Lettre {letter_type} ATI {ati.numero_ati}",
    )
    db.commit()

    return Response(
        content=buf.read(),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="lettre_{letter_type}_{ati.numero_ati}.pdf"'},
    )


@router.get("/exports/province/{province}/report.pdf")
async def export_province_report(
    province: str,
    current_user: User = Depends(require_roles(Role.admin, Role.ministre, Role.directeur)),
    db: Session = Depends(get_db),
):
    """Generate a PDF report for a specific province."""
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import HRFlowable, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    all_atis = db.execute(select(AgrementTechniqueIndustrielORM)).scalars().all()
    all_ops = (
        db.execute(select(OperateurIndustrielORM).where(OperateurIndustrielORM.is_active.is_(True))).scalars().all()
    )
    all_insp = db.execute(select(InspectionConformiteORM)).scalars().all()

    prov_ops = [o for o in all_ops if o.province == province]
    prov_atis = [a for a in all_atis if a.operateur and a.operateur.province == province]
    prov_insp = [i for i in all_insp if i.operateur and i.operateur.province == province]

    approved = sum(1 for a in prov_atis if a.statut == "approuve")
    conformes = sum(1 for i in prov_insp if i.statut_conformite == "conforme")

    now = now_utc()

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4, leftMargin=2.5 * cm, rightMargin=2.5 * cm, topMargin=2 * cm, bottomMargin=2 * cm
    )
    styles = getSampleStyleSheet()
    center = ParagraphStyle("center", parent=styles["Normal"], alignment=TA_CENTER)
    body = ParagraphStyle("body", parent=styles["Normal"], fontSize=11, leading=16)

    BLEU = colors.HexColor("#003F8F")
    VERT = colors.HexColor("#006233")

    story = []

    prov_label = province.replace("_", " ").title()

    story.append(
        Paragraph("REPUBLIQUE GABONAISE", ParagraphStyle("rg", parent=center, fontSize=11, textColor=colors.gray))
    )
    story.append(
        Paragraph(
            "Ministere de l'Industrie et de la Transformation Locale",
            ParagraphStyle("min", parent=center, fontSize=9, textColor=BLEU, fontName="Helvetica-Bold"),
        )
    )
    story.append(Spacer(1, 0.3 * cm))
    story.append(HRFlowable(width="100%", thickness=3, color=VERT))
    story.append(Spacer(1, 0.5 * cm))
    story.append(
        Paragraph(
            f"RAPPORT PROVINCIAL · {prov_label.upper()}",
            ParagraphStyle("title", parent=center, fontSize=18, textColor=BLEU, fontName="Helvetica-Bold"),
        )
    )
    story.append(
        Paragraph(
            f"Genere le {now.strftime('%d/%m/%Y')}",
            ParagraphStyle("date", parent=center, fontSize=10, textColor=colors.gray),
        )
    )
    story.append(Spacer(1, 0.6 * cm))

    # KPIs table
    kpis = [
        ["Indicateur", "Valeur"],
        ["Operateurs actifs", str(len(prov_ops))],
        ["ATIs soumis", str(len(prov_atis))],
        ["ATIs approuves", str(approved)],
        [
            "Taux d'approbation",
            f"{round(approved / max(len([a for a in prov_atis if a.statut in ('approuve', 'rejete')]), 1) * 100)}%",
        ],
        ["Inspections", str(len(prov_insp))],
        ["Taux de conformite", f"{round(conformes / max(len(prov_insp), 1) * 100)}%"],
    ]

    t = Table(kpis, colWidths=[8 * cm, 6 * cm])
    t.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), BLEU),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("PADDING", (0, 0), (-1, -1), 8),
                ("LINEBELOW", (0, 0), (-1, -2), 0.5, colors.HexColor("#e5e7eb")),
                ("LINEBELOW", (0, -1), (-1, -1), 1, BLEU),
            ]
        )
    )
    story.append(t)
    story.append(Spacer(1, 0.6 * cm))

    # Sector breakdown
    from collections import Counter

    sectors = Counter(a.secteur for a in prov_atis)
    if sectors:
        story.append(Paragraph("<b>Repartition sectorielle</b>", ParagraphStyle("h3", parent=body, textColor=BLEU)))
        story.append(Spacer(1, 0.2 * cm))
        for sect, count in sectors.most_common():
            story.append(Paragraph(f"• {sect.capitalize()} : {count} ATI(s)", body))

    story.append(Spacer(1, 1 * cm))
    story.append(HRFlowable(width="100%", thickness=1, color=VERT))
    story.append(
        Paragraph(
            "Document genere automatiquement par la PNPI",
            ParagraphStyle("footer", parent=center, fontSize=8, textColor=colors.gray, spaceBefore=8),
        )
    )

    doc.build(story)
    buf.seek(0)

    write_audit_event(
        db,
        actor=current_user.username,
        action="export.province_pdf",
        target=province,
        details=f"Rapport provincial {prov_label}",
    )
    db.commit()

    return Response(
        content=buf.read(),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="rapport_{province}.pdf"'},
    )


@router.post("/exports/create-archive")
async def create_digital_archive(
    data: dict,
    current_user: User = Depends(require_roles(Role.admin)),
    db: Session = Depends(get_db),
):
    """Create a sealed digital archive of ATI data for long-term preservation."""
    import hashlib
    import json as json_lib

    year = data.get("year", now_utc().year)

    all_atis = db.execute(select(AgrementTechniqueIndustrielORM)).scalars().all()
    year_atis = [a for a in all_atis if a.date_soumission.year == year]

    archive_data = {
        "metadata": {
            "type": "PNPI_DIGITAL_ARCHIVE",
            "version": "1.0",
            "year": year,
            "created_at": now_utc().isoformat(),
            "created_by": current_user.username,
            "record_count": len(year_atis),
        },
        "records": [
            {
                "numero_ati": a.numero_ati,
                "operateur": a.operateur.raison_sociale if a.operateur else None,
                "nif": a.operateur.nif if a.operateur else None,
                "secteur": a.secteur,
                "type_activite": a.type_activite,
                "statut": a.statut,
                "date_soumission": a.date_soumission.isoformat(),
                "date_decision": a.date_decision.isoformat() if a.date_decision else None,
                "sla_jours": a.sla_jours,
            }
            for a in year_atis
        ],
    }

    # Compute integrity hash
    content = json_lib.dumps(archive_data, sort_keys=True, ensure_ascii=False)
    integrity_hash = hashlib.sha256(content.encode()).hexdigest()
    archive_data["metadata"]["integrity_sha256"] = integrity_hash

    write_audit_event(
        db,
        actor=current_user.username,
        action="archive.create",
        target=str(year),
        details=f"Archive {year}: {len(year_atis)} ATIs, hash={integrity_hash[:16]}",
    )
    db.commit()

    return Response(
        content=json_lib.dumps(archive_data, indent=2, ensure_ascii=False),
        media_type="application/json",
        headers={
            "Content-Disposition": f'attachment; filename="archive_pnpi_{year}.json"',
            "X-PNPI-Integrity": integrity_hash,
        },
    )
